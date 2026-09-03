#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
fill_images.py — 轻风PPT 第4步（可选）：根据页面内容自动生成图片并插入图片占位符。

用法:
  python fill_images.py 生成结果.pptx
  python fill_images.py 生成结果.pptx --output 生成结果_生图.pptx
  python fill_images.py 生成结果.pptx --dry-run        # 只打印每图提示词，不调API、不改文件
  python fill_images.py 生成结果.pptx --size 2K
  python fill_images.py 生成结果.pptx --style "水彩插画风格，明亮温暖"
  python fill_images.py 生成结果.pptx --generator agnes   # 默认 agnes-image 2.1 flash

说明:
  - 自动检测所有 PICTURE 占位符；每个图按【位置最近】匹配一个关键词框，
    用「页标题 + 关键词 + 其正文」生成提示词（agnes 完美支持中文）。
  - 插入方式严格模拟"在占位符上手动插图"：把幻灯片侧空的 <p:sp type=pic> 占位符
    替换成 <p:pic>，并【深拷贝版式占位符的 spPr】（椭圆/圆角/旋转/位置等几何样式
    全部原样保留），仅把 <a:blip> 指向生成的图片。因此三栏的圆形、三图的弧形排布
    等版式样式都不会被破坏。
  - 按占位符宽高比选择 agnes 支持的 ratio；图片按「提示词+尺寸+模型」哈希缓存，
    重跑免费且秒回。
  - 单张生成失败不影响整体（该图留空 + 告警）。
  - 出图引擎默认 agnes-image 2.1 flash；不同用户 AGENT 环境可改 --generator
    （见 GENERATORS 字典，新增引擎只需加一个 call_xxx 函数）。
"""
import argparse, copy, hashlib, os, re, subprocess, sys, time
from lxml import etree
from pptx import Presentation

try:
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PPT
except Exception:
    PPT = None

AGNES_SCRIPT = os.path.join(
    os.path.expanduser('~'), '.workbuddy', 'skills', 'agnes-image',
    'scripts', 'generate_image.py')
MODEL = 'agnes-image-2.1-flash'
SUPPORTED_RATIOS = ['1:1', '16:9', '9:16', '4:3', '3:4', '2:3', '3:2', '21:9']
DEFAULT_STYLE = ('扁平插画风格，深蓝与青色科技感配色，简洁干净，留白充足，光影柔和，'
                 '主体突出；画面中不要出现任何文字、字母、数字或水印。')

EMU_PER_PX = 9525  # 914400 EMU/inch ÷ 96 px/inch

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def _A(tag):
    return '{%s}%s' % (A_NS, tag)


def _P(tag):
    return '{%s}%s' % (P_NS, tag)


def _R(attr):
    return '{%s}%s' % (R_NS, attr)


def is_picture(ph):
    try:
        if PPT is not None and ph.placeholder_format.type == PPT.PICTURE:
            return True
    except Exception:
        pass
    return 'PICTURE' in str(ph.placeholder_format.type)


def ph_type(ph):
    try:
        return str(ph.placeholder_format.type)
    except Exception:
        return ''


def nearest_ratio(w, h):
    if h <= 0:
        h = 1
    ar = w / h
    best, best_d = '1:1', 1e9
    for r in SUPPORTED_RATIOS:
        a, b = r.split(':')
        rr = int(a) / int(b)
        d = abs(ar - rr)
        if d < best_d:
            best_d, best = d, r
    return best


def pick_size(w, h):
    px = max(w, h) / EMU_PER_PX
    return '2K' if px >= 1200 else '1K'


def layout_kw_idx(layout):
    kw = set()
    for lph in layout.placeholders:
        if 'BODY' in ph_type(lph) and lph.text_frame.text.strip() == '输入关键词':
            kw.add(lph.placeholder_format.idx)
    return kw


def slide_title(slide):
    for ph in slide.placeholders:
        t = ph_type(ph)
        if 'TITLE' in t and 'SUBTITLE' not in t:
            return ph.text_frame.text.strip()
    for ph in slide.placeholders:
        if 'SUBTITLE' in ph_type(ph):
            return ph.text_frame.text.strip()
    return ''


def _norm(s):
    # 折叠换行/多空格为逗号，并清理因此产生的重复/相邻标点
    s = re.sub(r'\s+', '，', s.strip())
    s = re.sub(r'([，。])\1+', r'\1', s)
    s = re.sub(r'，。', '。', s)
    s = re.sub(r'。，', '。', s)
    return s


def build_prompt(title, keyword, body, style):
    core = title or ''
    if keyword:
        core = (core + '；' if core else '') + _norm(keyword)
    if body:
        body = _norm(body)
        if len(body) > 80:
            body = body[:80] + '…'
        core = (core + '：' if core else '') + body
    core = core.strip('；：，。 ')
    if not core:
        core = '教育培训主题插图'
    return core + '。' + style


def call_agnes(prompt, ratio, size, out_path):
    if not os.path.exists(AGNES_SCRIPT):
        raise RuntimeError('找不到 agnes-image 脚本: %s' % AGNES_SCRIPT)
    cmd = [sys.executable, AGNES_SCRIPT, prompt, '-o', out_path,
           '--size', size, '--ratio', ratio]
    r = subprocess.run(cmd, capture_output=True, text=True, encoding='utf-8')
    if r.returncode != 0 or not (os.path.exists(out_path) and os.path.getsize(out_path) > 0):
        raise RuntimeError('agnes 生图失败(%s): %s'
                           % (r.returncode, (r.stderr or r.stdout).strip()[:300]))
    return out_path


GENERATORS = {'agnes': call_agnes}


def generate_image(prompt, ratio, size, cache_dir, generator):
    os.makedirs(cache_dir, exist_ok=True)
    key = hashlib.sha256(('%s|%s|%s|%s' % (prompt, ratio, size, MODEL))
                         .encode('utf-8')).hexdigest()[:16]
    cache_file = os.path.join(cache_dir, key + '.png')
    if os.path.exists(cache_file) and os.path.getsize(cache_file) > 0:
        return cache_file, True
    GEN = GENERATORS.get(generator)
    if GEN is None:
        raise RuntimeError('未知出图引擎: %s（可选: %s）'
                           % (generator, ','.join(GENERATORS)))
    last = None
    for attempt in range(3):           # 应对限流(429)等瞬时失败
        try:
            GEN(prompt, ratio, size, cache_file)
            return cache_file, False
        except Exception as e:
            last = e
            if attempt < 2:
                time.sleep(3 * (attempt + 1))
                continue
    raise last


def fill_placeholder(slide, ph_shape, layout_ph_shape, img_path):
    """把幻灯片侧空的 <p:sp type=pic> 占位符替换成携图的 <p:pic>，
    深拷贝版式占位符的 spPr（椭圆/圆角/旋转/位置等几何样式全部保留），
    等价于在占位符上手动插入图片。返回新图片的 rId。"""
    image_part, rId = slide.part.get_or_add_image_part(img_path)
    ph_el = ph_shape.element

    pic = etree.Element(_P('pic'))

    # nvPicPr
    nvPicPr = etree.SubElement(pic, _P('nvPicPr'))
    orig_cNvPr = ph_el.find(_P('nvSpPr') + '/' + _P('cNvPr'))
    cNvPr = etree.SubElement(nvPicPr, _P('cNvPr'))
    cNvPr.set('id', orig_cNvPr.get('id') if orig_cNvPr is not None else '100')
    cNvPr.set('name', (orig_cNvPr.get('name') + ' (filled)')
              if orig_cNvPr is not None and orig_cNvPr.get('name') else 'Picture')
    etree.SubElement(nvPicPr, _P('cNvPicPr'))
    nvPr = etree.SubElement(nvPicPr, _P('nvPr'))
    orig_ph = ph_el.find(_P('nvSpPr') + '/' + _P('nvPr') + '/' + _P('ph'))
    if orig_ph is not None:
        nvPr.append(copy.deepcopy(orig_ph))

    # blipFill —— 指向新图，以 fillRect 拉伸（"填充"行为，等同手动插图）
    blipFill = etree.SubElement(pic, _P('blipFill'))
    blip = etree.SubElement(blipFill, _A('blip'))
    blip.set(_R('embed'), rId)
    stretch = etree.SubElement(blipFill, _A('stretch'))
    etree.SubElement(stretch, _A('fillRect'))

    # spPr —— 直接深拷贝版式占位符的 spPr（几何/位置/旋转/特效全部原样保留）
    if layout_ph_shape is not None:
        src_spPr = layout_ph_shape.element.find(_P('spPr'))
        if src_spPr is not None:
            pic.append(copy.deepcopy(src_spPr))
    if pic.find(_P('spPr')) is None:
        etree.SubElement(pic, _P('spPr'))

    # 用新 <p:pic> 替换原空占位符
    parent = ph_el.getparent()
    parent.replace(ph_el, pic)
    return rId


def main():
    ap = argparse.ArgumentParser(description='轻风PPT 第4步：按内容自动生图插入图片占位符')
    ap.add_argument('pptx', help='第3步生成的 pptx 路径')
    ap.add_argument('--output', help='输出路径（默认 <名>_生图.pptx）')
    ap.add_argument('--cache', help='图片缓存目录（默认与 pptx 同目录 <名>_imgcache）')
    ap.add_argument('--size', default=None, help='尺寸档位 1K/2K（默认按占位符大小自动）')
    ap.add_argument('--style', default=DEFAULT_STYLE, help='出图风格后缀')
    ap.add_argument('--generator', default='agnes', help='出图引擎（默认 agnes）')
    ap.add_argument('--dry-run', action='store_true',
                    help='只打印每图提示词，不调API、不改文件')
    args = ap.parse_args()

    if not os.path.exists(args.pptx):
        print('[错误] 文件不存在:', args.pptx)
        sys.exit(1)

    stem, ext = os.path.splitext(args.pptx)
    out_path = args.output or (stem + '_生图' + ext)
    cache_dir = args.cache or (stem + '_imgcache')

    prs = Presentation(args.pptx)
    plan = []

    for si, slide in enumerate(prs.slides, 1):
        lay = slide.slide_layout
        kw_idx = layout_kw_idx(lay)
        pics = [ph for ph in slide.placeholders if is_picture(ph)]
        if not pics:
            continue
        # 版式同名图片占位符（按 idx 精确匹配），用于深拷贝几何样式
        lmap = {}
        for lph in lay.placeholders:
            if is_picture(lph):
                try:
                    lmap[lph.placeholder_format.idx] = lph
                except Exception:
                    pass
        kw_phs = [ph for ph in slide.placeholders
                  if ph.placeholder_format.idx in kw_idx and ph.text_frame.text.strip()]
        content_phs = [ph for ph in slide.placeholders
                       if 'BODY' in ph_type(ph)
                       and ph.placeholder_format.idx not in kw_idx
                       and ph.text_frame.text.strip()]
        title = slide_title(slide)
        # 关键词 → 正文 按顺序配对
        kw_body = {}
        for i, kp in enumerate(kw_phs):
            bd = content_phs[i].text_frame.text if i < len(content_phs) else ''
            kw_body[kp.placeholder_format.idx] = (kp.text_frame.text.strip(), bd.strip())
        pics_sorted = sorted(pics, key=lambda p: (p.top, p.left))
        kw_sorted = sorted(kw_phs, key=lambda p: (p.top, p.left))
        # 按位置排序后一一对应：第 i 个图 ↔ 第 i 个关键词（网格布局的正确语义）
        for i, pic in enumerate(pics_sorted):
            keyword, body = ('', '')
            if i < len(kw_sorted):
                kp = kw_sorted[i]
                keyword, body = kw_body.get(kp.placeholder_format.idx, ('', ''))
            layout_ph = lmap.get(pic.placeholder_format.idx)
            if layout_ph is not None:
                g = layout_ph
                ratio = nearest_ratio(g.width, g.height)
                size = args.size or pick_size(g.width, g.height)
            else:
                g = pic
                ratio = nearest_ratio(pic.width or 1, pic.height or 1)
                size = args.size or pick_size(pic.width or 1, pic.height or 1)
            prompt = build_prompt(title, keyword, body, args.style)
            plan.append((si, lay.name, pic.placeholder_format.idx, keyword,
                         ratio, size, prompt, slide, pic, layout_ph))

    if args.dry_run:
        print('=== DRY-RUN：以下为将生成的图片计划（不调API、不改文件）===')
        for si, ln, idx, kw, ratio, size, prompt, _ in plan:
            print('  第%d页 [%s] 图框#%s  ratio=%s size=%s' % (si, ln, idx, ratio, size))
            print('    关键词: %s' % (kw or '(无)'))
            print('    提示词: %s' % prompt)
        print('共 %d 张图待生成。' % len(plan))
        return

    total = gen = reused = failed = 0
    for si, ln, idx, kw, ratio, size, prompt, slide, pic, layout_ph in plan:
        total += 1
        try:
            img, was_cached = generate_image(prompt, ratio, size, cache_dir, args.generator)
            # 保留占位符几何样式（椭圆/弧形/位置），仅替换图片内容
            fill_placeholder(slide, pic, layout_ph, img)
            if was_cached:
                reused += 1
                print('  第%d页 图框#%s: 复用缓存 %s' % (si, idx, os.path.basename(img)))
            else:
                gen += 1
                print('  第%d页 图框#%s: 已生成并插入 (%s, %s)' % (si, idx, ratio, size))
        except Exception as e:
            failed += 1
            print('  ⚠ 第%d页 图框#%s 生图失败，留空: %s' % (si, idx, e))

    prs.save(out_path)
    print('\n完成: 共 %d 张图 | 新生成 %d | 复用缓存 %d | 失败 %d'
          % (total, gen, reused, failed))
    print('输出: %s' % out_path)
    print('缓存目录: %s' % cache_dir)


if __name__ == '__main__':
    main()
