# -*- coding: utf-8 -*-
"""qingfeng-ppt · 第1步：把用户上传的 PPT 模板优化为「标准母版版式」

规则（与技能定义一致）：
  1. 所有文本占位符默认文字颜色 -> 白色 (FFFFFF)，保证深色/蓝色背景图上可见
  2. 文本框固定大小：去掉随文字变形的 spAutoFit 与禁止自动调整的 noAutofit，
     统一用 normAutofit —— 即【根据文本框自动调整字号】，超框自动缩字、排版不乱
  3. 标题占位符 wrap=none -> 强制单行；正文占位符 wrap=square -> 允许换行
  4. 标题与正文都保留 normAutofit：文字多/长时自动缩字号，永不出框、不破版

铁律（不可违反）：
  - 只改【文本属性】（颜色、自动调整），绝不改动背景图、整体配色主题、图形、
    母版/版式的数量与结构 —— 所以“背景/风格与用户模板一致”，只是把文字规则
    标准化，方便后面纯灌入。
  - 用户母版的这 14 个版式（封面/目录/节标题/单栏/两栏/三栏/左图右文/右图左文/
    大图/四图/三图/四栏/六栏/空白）原样保留，一个都不增删。

用法:
    python optimize_template.py <输入模板.pptx> [输出模板.pptx]
不写输出路径时，默认在原文件名后加 "_optimized"。
"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
WHITE = 'FFFFFF'
# 占位符类型名前缀（来自 str(ph.placeholder_format.type) 的 "TITLE (1)" 形式）
TITLE_TYPES = {'TITLE', 'CENTER_TITLE', 'VERTICAL_TITLE'}


def _q(tag):
    return qn(tag)


def set_white_defrpr(pPr):
    """在 pPr 下确保 defRPr/solidFill/srgbClr=白色"""
    defRPr = pPr.find(_q('a:defRPr'))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, _q('a:defRPr'))
    for sf in defRPr.findall(_q('a:solidFill')):
        defRPr.remove(sf)
    sf = etree.SubElement(defRPr, _q('a:solidFill'))
    srgb = etree.SubElement(sf, _q('a:srgbClr'))
    srgb.set('val', WHITE)
    # 清掉可能存在的其他颜色节点（如 sysClr）
    for c in defRPr.findall(_q('a:sysClr')):
        defRPr.remove(c)


def white_out_styles(txstyles_elem):
    """遍历 a:titleStyle/bodyStyle/otherStyle 下的 defPPr / lvlNpPr，置白"""
    if txstyles_elem is None:
        return
    for pPr in txstyles_elem.iter(_q('a:defPPr')):
        set_white_defrpr(pPr)
    for i in range(1, 10):
        for pPr in txstyles_elem.iter(_q(f'a:lvl{i}pPr')):
            set_white_defrpr(pPr)


def white_out_lststyle(txBody):
    """遍历占位符 txBody 的 lstStyle 各级，置白（若无则创建）"""
    lst = txBody.find(_q('a:lstStyle'))
    if lst is None:
        lst = etree.SubElement(txBody, _q('a:lstStyle'))
    for i in range(1, 10):
        tag = _q(f'a:lvl{i}pPr')
        pPr = lst.find(tag)
        if pPr is None:
            pPr = etree.SubElement(lst, tag)
        set_white_defrpr(pPr)
    # 同时处理没有层级时的默认 defPPr
    defP = lst.find(_q('a:defPPr'))
    if defP is None:
        defP = etree.SubElement(lst, _q('a:defPPr'))
    set_white_defrpr(defP)


def fix_bodyPr(bodyPr, is_title):
    if bodyPr is None:
        return
    # 去掉随文字变形的 spAutoFit（按文字多少变形，排版会乱）
    for sa in bodyPr.findall(_q('a:spAutoFit')):
        bodyPr.remove(sa)
    # 去掉“禁止自动调整”，确保【根据文本框自动调整字号】被保留
    for na in bodyPr.findall(_q('a:noAutofit')):
        bodyPr.remove(na)
    # 统一 normAutofit（超框自动缩字号，排版不乱）—— 已存在则保留不动
    if bodyPr.find(_q('a:normAutofit')) is None:
        etree.SubElement(bodyPr, _q('a:normAutofit'))
    # 标题强制单行，正文允许换行
    bodyPr.set('wrap', 'none' if is_title else 'square')


def optimize_presentation(prs):
    # 1) 母版 txStyles 置白
    master = prs.slide_masters[0]
    master_element = master.element
    txStyles = master_element.find(_q('p:txStyles'))
    white_out_styles(txStyles)

    # 母版自带占位符（如标题/页脚）也标准化
    for ph in master.placeholders:
        _optimize_placeholder(ph)

    # 2) 每个版式的占位符标准化
    for lay in prs.slide_masters[0].slide_layouts:
        for ph in lay.placeholders:
            _optimize_placeholder(ph)


def _optimize_placeholder(ph):
    txBody = ph.element.find(_q('p:txBody'))
    if txBody is None:
        return
    is_title = str(ph.placeholder_format.type).split(' ')[0] in TITLE_TYPES
    white_out_lststyle(txBody)
    bodyPr = txBody.find(_q('a:bodyPr'))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, _q('a:bodyPr'))
    fix_bodyPr(bodyPr, is_title)


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    src = sys.argv[1]
    if len(sys.argv) >= 3:
        dst = sys.argv[2]
    else:
        base, ext = os.path.splitext(src)
        dst = base + '_optimized' + ext

    prs = Presentation(src)
    optimize_presentation(prs)
    prs.save(dst)
    print('已优化 ->', dst)


if __name__ == '__main__':
    main()
