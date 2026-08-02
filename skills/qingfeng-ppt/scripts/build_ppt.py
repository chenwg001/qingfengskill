# -*- coding: utf-8 -*-
"""qingfeng-ppt · 第3步：按大纲把文字【纯灌入】模板文本框

设计铁律（不可违反）:
  - 只用模板原生【文本】占位符灌字，绝不新建文本框、绝不改占位符的位置/大小/
    颜色/字号/换行等任何样式，不添加任何元素。
  - 颜色、字号、是否单行等全部由模板决定；生成后若不满意，由用户在 PowerPoint
    里自己调。本脚本在生成阶段一律不碰。
  - 图片 / 其他非文本占位符一律【留空】，等用户手动添加。
  - 版式由大纲页标题末尾的 [版式] 标注决定；无标注时按二级标题(关键词)数量自动匹配。

大纲（md）格式（严格）:
    # 页标题                 ← 一级标题 → 标题文本框
    ## 关键词一               ← 二级标题 → 关键词文本框
    ## 关键词二
    正文段落一                ← 普通段落 → 添加文本的文本框
    正文段落二
    # 另一页 [两栏]           ← [版式] 标注可锁定该页版式

特殊版式:
    [封面]   标题→标题框，正文第一段→副标题框
    [目录]   标题→标题框，正文逐段→正文框
    [节标题] 标题→标题框
    [空白]   生成干净空页

用法:
    python build_ppt.py 大纲.md
    python build_ppt.py pages.json        # 兼容旧式 JSON（见下文）
可选环境变量 TEMPLATE 指定模板；否则用技能自带默认模板。

pages.json（可选/兼容）:
    {"template":"...","output":"...","pages":[
       {"layout":"cover","title":"主标题","subtitle":"副标题"},
       {"title":"页标题","keywords":["关键词一","关键词二"],"body":["正文一","正文二"]}
    ]}
"""
import sys, io, os, re, json
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_TEMPLATE = os.path.join(HERE, '..', 'assets', '轻风模板_规范版.pptx')

# 版式 key → 名称关键字（用于把标注/自动选择映射到模板真实版式）
NAME_MAP = [
    ('cover',     ['封面', 'cover']),
    ('toc',       ['目录', 'toc']),
    ('section',   ['节', 'section']),
    ('1col',      ['单栏', '1col']),
    ('2col',      ['两栏', '2col']),
    ('3col',      ['三栏', '3col']),
    ('4col',      ['四栏', '4col']),
    ('6col',      ['六栏', '6col']),
    ('img_left',  ['左图', 'imgleft']),
    ('img_right', ['右图', 'imgright']),
    ('full_img',  ['大图', 'imagefull', 'full']),
    ('4images',   ['四图', '四宫格', '4images']),
    ('3images',   ['三图', '3images']),
    ('blank',     ['空白', 'blank']),
]
ALIASES = {
    '封面': 'cover', '目录': 'toc', '节标题': 'section', '章节': 'section',
    '单栏': '1col', '两栏': '2col', '三栏': '3col', '四栏': '4col', '六栏': '6col',
    '左图右文': 'img_left', '右图左文': 'img_right', '左文右图': 'img_right',
    '大图': 'full_img', '全图': 'full_img',
    '四图': '4images', '三图': '3images', '空白': 'blank',
    '00': 'cover', '01': 'toc', '02': 'section', '03': '1col', '04': '2col',
    '05': '3col', '06': 'img_left', '07': 'img_right', '08': 'full_img',
    '09': '4images', '10': '3images', '11': 'blank',
    '12': '4col', '13': '6col',
}


def resolve_layout_map(prs):
    """key -> 模板版式索引，按名称关键字匹配；未命中按序兜底"""
    layouts = prs.slide_masters[0].slide_layouts
    names = [lay.name.lower() if lay.name else '' for lay in layouts]
    used = set()
    mapping = {}
    for key, kws in NAME_MAP:
        for i, nm in enumerate(names):
            if i in used:
                continue
            if any(kw in nm for kw in kws):
                mapping[key] = i
                used.add(i)
                break
    for key, _, in NAME_MAP:
        if key in mapping:
            continue
        for i in range(len(layouts)):
            if i not in used:
                mapping[key] = i
                used.add(i)
                break
    return mapping


def norm_layout(name):
    if not name:
        return None
    key = str(name).strip().lower()
    key = ALIASES.get(key, ALIASES.get(str(name).strip(), key))
    return key if key in dict(NAME_MAP) else None


def auto_layout(page):
    """无 [版式] 标注时，按二级标题(关键词)数量自动选版式"""
    n_kw = len(page.get('keywords') or [])
    n_body = len(page.get('body') or [])
    if n_kw == 0 and n_body == 0:
        return 'section' if page.get('title') else 'blank'
    if n_kw == 0:
        return '1col'
    return {1: '2col', 2: '2col', 3: '3col', 4: '4col'}.get(n_kw, '6col')


def classify(slide):
    """把一页的占位符分成 标题/副标题/关键词框/正文框，不做任何修改。

    注意：幻灯片实例的占位符文本是空的，'输入关键词' 提示文字只存在于【版式
    (layout)定义】里。因此关键词框的判定要先在 layout 上按 idx 识别，再映射回
    幻灯片占位符。
    """
    layout = slide.slide_layout
    kw_idx = set()
    for lph in layout.placeholders:
        try:
            t = str(lph.placeholder_format.type)
        except Exception:
            t = ''
        if 'BODY' in t and lph.text_frame.text.strip() == '输入关键词':
            kw_idx.add(lph.placeholder_format.idx)

    title = subtitle = None
    kw, content = [], []
    for ph in slide.placeholders:
        t = str(ph.placeholder_format.type)
        idx = ph.placeholder_format.idx
        if 'SUBTITLE' in t:          # 必须放在 TITLE 之前：'SUBTITLE' 含 'TITLE' 子串
            subtitle = ph
        elif 'TITLE' in t:
            title = ph
        elif 'BODY' in t:
            if idx in kw_idx:
                kw.append(ph)
            else:
                content.append((idx, ph))
    content.sort(key=lambda x: x[0])
    return {'title': title, 'subtitle': subtitle,
            'kw': kw, 'content': [p for _, p in content]}


def _set(ph, text):
    """纯灌入：只写文字，不改任何样式"""
    tf = ph.text_frame
    lines = str(text).split('\n')
    tf.text = lines[0] if lines else ''
    for ln in lines[1:]:
        tf.add_paragraph().text = ln


def _clear(ph):
    ph.text_frame.text = ''


def pour(slide, page, layout_key):
    parts = classify(slide)
    warnings = []
    title = page.get('title') or ''

    # 标题
    if parts['title'] is not None and title:
        _set(parts['title'], title)

    if layout_key == 'cover':
        body = page.get('body') or []
        if parts['subtitle'] is not None and body:
            _set(parts['subtitle'], body[0])
        return warnings

    if layout_key == 'section':
        return warnings  # 仅标题

    if layout_key == 'blank':
        return warnings  # 干净空页，什么都不填

    if layout_key == 'toc':
        body = page.get('body') or []
        if parts['content']:
            box = parts['content'][0]
            _set(box, body[0] if body else '')
            for ln in body[1:]:
                box.text_frame.add_paragraph().text = ln
        return warnings

    # 通用内容版式（单栏/两栏/三栏/四栏/六栏/图文/图集）
    keywords = page.get('keywords') or []
    body = page.get('body') or []

    # 关键词 → 关键词框
    for i, kwph in enumerate(parts['kw']):
        if i < len(keywords):
            _set(kwph, keywords[i])
        else:
            _clear(kwph)  # 没用到的关键词框清掉"输入关键词"提示

    # 正文 → 正文框
    cb = parts['content']
    if cb:
        for i, cph in enumerate(cb):
            if i < len(body):
                _set(cph, body[i])
            # 没用到的正文框留空（模板默认即空）
        # 正文段数超过正文框数：多出的段追加到最后一只框
        if len(body) > len(cb):
            extra = body[len(cb):]
            last = cb[-1].text_frame
            for ln in extra:
                last.add_paragraph().text = ln
            warnings.append(
                f'正文段数 {len(body)} 超过正文框 {len(cb)}，多余已并入最后一只框')
    elif body:
        warnings.append('该版式无可用正文框，正文未填入')

    if keywords and not parts['kw']:
        warnings.append('该版式无关键词框，二级标题未填入')
    return warnings


def parse_outline(md_path):
    """解析严格 md 大纲 -> pages 列表。

    规则（与 SKILL 第 2 步方法论一致）：
      # 页标题            -> 新页；一级标题 -> 标题文本框
      ## 关键词           -> 该页一个关键词（二级标题）-> 关键词文本框
      普通段落            -> 紧跟在哪个 ## 之后，就属于那个关键词的【核心内容】
                             （实现「关键词」与「核心内容」一一配对）
      # 页标题 之下、首个 ## 之前的普通段落 -> 页前引言，并入第一个关键词内容
      （封面副标题 / 目录项 / 单栏引言等无关键词的页，正文整体进入 body）

    返回 page: {title, keywords:[...], body:[...], layout}
      - 有关键词的页：len(keywords) == len(body)，body[i] 是 keywords[i] 的核心内容
      - 无关键词的页：keywords=[]，body=[整段文本...]
    """
    pages, cur = [], None
    cur_kw = -1  # 当前关键词在 page['keywords'] 中的下标
    with open(md_path, encoding='utf-8') as f:
        for raw in f:
            s = raw.strip()
            if s.startswith('# '):
                cur = {'title': '', 'keywords': [], 'body': [], 'layout': None, '_pre': []}
                pages.append(cur)
                cur_kw = -1
                rest = s[2:].strip()
                m = re.search(r'\[([^\]]+)\]\s*$', rest)
                if m:
                    cur['layout'] = m.group(1).strip()
                    rest = rest[:m.start()].strip()
                cur['title'] = rest
            elif s.startswith('## '):
                if cur is not None:
                    cur['keywords'].append(s[3:].strip())
                    cur['body'].append('')   # 为该关键词预留核心内容槽位
                    cur_kw = len(cur['keywords']) - 1
            elif s == '':
                continue
            else:
                if cur is not None:
                    if cur_kw >= 0:
                        # 属于当前关键词的核心内容
                        cur['body'][cur_kw] = (cur['body'][cur_kw] + '\n' + s) if cur['body'][cur_kw] else s
                    else:
                        # 页标题下、首个关键词之前的引言
                        cur['_pre'].append(s)
    # 收尾：把 _pre（页前引言）归位
    for p in pages:
        pre = p.pop('_pre', [])
        if not pre:
            continue
        if p['keywords'] and p['body']:
            # 并入第一个关键词的核心内容（保留本就有的内容）
            lead = '\n'.join(pre)
            p['body'][0] = (lead + '\n' + p['body'][0]) if p['body'][0] else lead
        else:
            # 无关键词页（封面/目录/单栏）：正文整体
            p['body'] = pre + p['body']
    return [p for p in pages if p.get('title') or p.get('keywords') or p.get('body')]


def build(cfg):
    template = cfg.get('template') or os.environ.get('TEMPLATE') or DEFAULT_TEMPLATE
    if not os.path.exists(template):
        print('[错误] 模板不存在:', template)
        sys.exit(1)
    output = cfg['output']
    prs = Presentation(template)

    # 清空模板自带示例页（仅为移除样板，不改任何样式）
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        rId = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)

    layout_map = resolve_layout_map(prs)
    # 模板可能缺某版式（如用户手动删掉「单栏」），缺失时优雅降级到最近的真实版式，不崩溃
    FALLBACK = {
        '1col': '2col', '2col': '3col', '3col': '2col', '4col': '3col',
        '6col': '4col', 'img_left': '2col', 'img_right': '2col',
        'full_img': 'blank', '4images': 'blank', '3images': 'blank',
    }
    all_warnings = []
    for i, page in enumerate(cfg['pages'], 1):
        key = norm_layout(page.get('layout')) or auto_layout(page)
        if key not in layout_map:
            fb = FALLBACK.get(key, 'blank')
            if fb not in layout_map:
                fb = next(iter(layout_map))
            all_warnings.append(f'第{i}页: 版式 {key} 模板无对应，已回退为 {fb}')
            key = fb
        slide = prs.slides.add_slide(prs.slide_layouts[layout_map[key]])
        ws = pour(slide, page, key)
        all_warnings.extend([f'第{i}页: {w}' for w in ws])
        print(f'第{i:02d}页  版式={key}  标题={page.get("title","")!r}'
              f'  关键词={len(page.get("keywords") or [])}  正文={len(page.get("body") or [])}')

    prs.save(output)
    print('已生成 ->', output)
    for w in all_warnings:
        print('[警告]', w)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    src = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else \
        os.path.splitext(src)[0] + '_生成.pptx'
    if src.lower().endswith('.md'):
        pages = parse_outline(src)
        cfg = {'output': output, 'pages': pages}
    else:
        cfg = json.load(open(src, encoding='utf-8'))
        cfg.setdefault('output', output)
    build(cfg)
