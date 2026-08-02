# -*- coding: utf-8 -*-
"""qingfeng-ppt · 第1步：把用户模板的「母版背景图 + 主题配色」换到轻风 14 版式母版骨架上

设计原则（不可违反）:
  - 只动【视觉皮肤】：母版背景图、主题配色；绝不改版式结构、占位符位置/数量、
    关键词框、分栏、文字框大小等任何布局。
  - 换肤后按【背景亮度自动选黑/白字】（luminance 判定），不再强制白字。
  - 内置轻风模板【不改】时保持白字（其默认即白字，本脚本不碰）。
  - 只读【母版】背景；母版无背景 → 打印提示让用户上传图片，退出。
  - 若上传模板本身规范（母版含 标题/关键词/正文 三级文本占位符），则直接使用
    上传模板，不换肤。

用法:
  # 上传了模板（非规范则换肤；规范则直接用）
  python restyle_template.py --user 用户模板.pptx --output 更新后模板.pptx
  # 直接给了一张背景图片
  python restyle_template.py --image 背景图.png --output 更新后模板.pptx
  # 模板 + 图片（图片优先作背景）
  python restyle_template.py --user 用户模板.pptx --image 背景图.png --output 更新后模板.pptx
  # 指定内置骨架（默认读技能自带 轻风模板_规范版.pptx）
  python restyle_template.py --user 用户模板.pptx --base 轻风模板_规范版.pptx --output 更新后模板.pptx
"""
import sys, os, io, copy, argparse, tempfile
sys.stdout = __import__('io').TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = __import__('io').TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from lxml import etree
from PIL import Image

# ---- 命名空间 ----
P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
THEME_RTYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'


def pq(tag):
    return '{%s}%s' % (P, tag)


def aq(tag):
    return '{%s}%s' % (A, tag)


def rq(tag):
    return '{%s}%s' % (R, tag)


# ---------------------------------------------------------------------------
# 1) 规范模板判定：母版含 标题 / 关键词 / 正文 三级文本占位符
#    关键签名 = 版式里存在「版式层文本 == 输入关键词」的 BODY 占位符（关键词框），
#    这正是 build_ppt.py 的 classify() 用来识别关键词框的依据；普通 Office 模板没有。
# ---------------------------------------------------------------------------
def is_standard_template(prs):
    m = prs.slide_masters[0]
    layouts = m.slide_layouts
    has_kw_box = False
    title_layouts = 0
    for lay in layouts:
        titles = [ph for ph in lay.placeholders
                  if 'TITLE' in str(ph.placeholder_format.type)]
        if titles:
            title_layouts += 1
        for ph in lay.placeholders:
            if ('BODY' in str(ph.placeholder_format.type)
                    and ph.text_frame.text.strip() == '输入关键词'):
                has_kw_box = True
    # 版式名需覆盖主要角色（封面/目录/节/分栏），进一步排除普通 Office 模板
    names = [(lay.name or '').lower() for lay in layouts]
    role_kw = {
        'cover':   ['封面', 'cover'],
        'toc':     ['目录', 'toc'],
        'section': ['节', 'section'],
        'col':     ['栏', 'col'],
    }
    covered = sum(1 for kws in role_kw.values()
                  if any(any(kw in nm for kw in kws) for nm in names))
    return has_kw_box and title_layouts >= 4 and covered >= 3


# ---------------------------------------------------------------------------
# 2) 提取母版背景：返回 ('image', bytes, ext) / ('solid', 'RRGGBB') / None
# ---------------------------------------------------------------------------
def get_master_bg(prs):
    m = prs.slide_masters[0]
    cSld = m.element.find(pq('cSld'))
    if cSld is None:
        return None
    bg = cSld.find(pq('bg'))
    if bg is None:
        return None
    bgPr = bg.find(pq('bgPr'))
    if bgPr is None:
        return None
    # 图片背景
    blipFill = bgPr.find(pq('blipFill'))
    if blipFill is not None:
        blip = blipFill.find(aq('blip'))
        if blip is not None:
            rid = blip.get(rq('embed'))
            if rid:
                part = m.part.related_part(rid)
                ext = os.path.splitext(part.partname)[1].lstrip('.') or 'png'
                return ('image', part.blob, ext)
    # 纯色背景
    solid = bgPr.find(pq('solidFill'))
    if solid is None:
        solid = bgPr.find(aq('solidFill'))
    if solid is not None:
        srgb = solid.find(aq('srgbClr'))
        if srgb is not None:
            return ('solid', srgb.get('val'))
    return None


# ---------------------------------------------------------------------------
# 3) 亮度判定：返回 'FFFFFF' 或 '000000'
# ---------------------------------------------------------------------------
def bg_luminance_color(kind, data):
    if kind == 'solid':
        hexv = data if isinstance(data, str) else ''
        if len(hexv) == 6:
            r, g, b = int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16)
        else:
            r = g = b = 128
    else:  # image bytes
        im = Image.open(io.BytesIO(data)).convert('RGB')
        w, h = (60, 60)
        im = im.resize((w, h))
        px = im.load()
        sr = sg = sb = 0
        for y in range(h):
            for x in range(w):
                r, g, b = px[x, y]
                sr += r; sg += g; sb += b
        n = w * h
        r, g, b = sr / n, sg / n, sb / n
    L = 0.299 * r + 0.587 * g + 0.114 * b   # 0-255
    return '000000' if L > 140 else 'FFFFFF'  # 浅底→黑字；深底→白字


# ---------------------------------------------------------------------------
# 4) 设置母版背景（图片或纯色）
# ---------------------------------------------------------------------------
def _build_bg_image(rid):
    xml = (
        '<p:bg xmlns:p="%s"><p:bgPr>'
        '<a:blipFill xmlns:a="%s"><a:blip xmlns:r="%s" r:embed="%s"/>'
        '<a:stretch><a:fillRect/></a:stretch></a:blipFill>'
        '</p:bgPr></p:bg>'
    ) % (P, A, R, rid)
    return etree.fromstring(xml.encode('utf-8'))


def _build_bg_solid(color):
    xml = (
        '<p:bg xmlns:p="%s"><p:bgPr>'
        '<a:solidFill xmlns:a="%s"><a:srgbClr val="%s"/></a:solidFill>'
        '</p:bgPr></p:bg>'
    ) % (P, A, color)
    return etree.fromstring(xml.encode('utf-8'))


def _replace_bg_picture(base_master, blob, ext, slide_w, slide_h):
    """若母版用 spTree 里的一张全屏图片作【可见背景】，则替换其图像内容
    （保留位置/尺寸、其上方的蒙版层与装饰）。返回是否成功。

    修复点：仅改 p:cSld/p:bg 时，原母版的全屏图片仍压在上层，肉眼看「换肤失败」。"""
    cSld = base_master.element.find(pq('cSld'))
    if cSld is None:
        return False
    spTree = cSld.find(pq('spTree'))
    if spTree is None:
        return False
    # 画面尺寸（EMU），用于判定「全屏」背景图（由调用方传入，避免依赖 package 内部属性）
    W, H = slide_w, slide_h
    if W == 0 or H == 0:
        W, H = 12192000, 6858000   # 兜底：默认 16:9
    pics = [ch for ch in spTree if etree.QName(ch).localname == 'pic']
    pics = [ch for ch in spTree if etree.QName(ch).localname == 'pic']
    target = None
    for pic in pics:
        xfrm = pic.find('.//' + aq('xfrm'))
        if xfrm is None:
            continue
        off = xfrm.find(aq('off'))
        ext_el = xfrm.find(aq('ext'))
        if off is None or ext_el is None:
            continue
        x = int(off.get('x', '0')); y = int(off.get('y', '0'))
        w = int(ext_el.get('cx', '0')); h = int(ext_el.get('cy', '0'))
        # 全屏背景：贴左上角且覆盖 >=85% 画面
        if x <= 1000 and y <= 1000 and w >= 0.85 * W and h >= 0.85 * H:
            target = pic
            break
    if target is None and pics:
        target = pics[0]          # 兜底：取最底部图片
    if target is None:
        return False
    tmp = tempfile.NamedTemporaryFile(suffix='.' + ext, delete=False)
    tmp.write(blob)
    tmp.close()
    rel = base_master.part.get_or_add_image_part(tmp.name)
    new_rid = rel[1] if isinstance(rel, tuple) else rel
    os.unlink(tmp.name)
    blip = target.find('.//' + aq('blip'))
    if blip is None:
        return False
    old_rid = blip.get(rq('embed'))
    blip.set(rq('embed'), new_rid)
    if old_rid and old_rid != new_rid:
        try:
            base_master.part.drop_rel(old_rid)
        except Exception:
            pass
    return True


def set_master_bg(base_master, kind, data, slide_w=12192000, slide_h=6858000):
    """在 base_master 母版上设置背景；返回 ('image',bytes) 或 ('solid',color) 供亮度判定。

    优先替换母版 spTree 里的全屏背景图片（可见背景层），保留其上的蒙版层/装饰；
    仅当模板无全屏图片层时，才回退到 p:cSld/p:bg。"""
    cSld = base_master.element.find(pq('cSld'))

    def _clear_old_bg():
        old = cSld.find(pq('bg'))
        if old is not None:
            blip = old.find('.//' + aq('blip'))
            if blip is not None:
                rid = blip.get(rq('embed'))
                if rid:
                    try:
                        base_master.part.drop_rel(rid)
                    except Exception:
                        pass
            cSld.remove(old)

    if kind == 'image':
        blob, ext = data
        if _replace_bg_picture(base_master, blob, ext, slide_w, slide_h):
            _clear_old_bg()       # 移除冗余 p:cSld/p:bg，避免两层背景叠加
            return ('image', blob)
        # 回退：模板无全屏图片层 → 改 p:cSld/p:bg
        _clear_old_bg()
        tmp = tempfile.NamedTemporaryFile(suffix='.' + ext, delete=False)
        tmp.write(blob)
        tmp.close()
        rel = base_master.part.get_or_add_image_part(tmp.name)
        rid = rel[1] if isinstance(rel, tuple) else rel
        os.unlink(tmp.name)
        bg_el = _build_bg_image(rid)
        cSld.insert(0, bg_el)
        return ('image', blob)
    else:
        color = data
        _clear_old_bg()
        bg_el = _build_bg_solid(color)
        cSld.insert(0, bg_el)
        return ('solid', color)


# ---------------------------------------------------------------------------
# 5) 复制主题配色（clrScheme）到 base 主题
# ---------------------------------------------------------------------------
def _theme_el(theme_part):
    """健壮取 theme 的 XML 元素：XmlPart 有 .element，普通 Part 退回解析 blob"""
    el = getattr(theme_part, 'element', None)
    if el is None:
        el = etree.fromstring(theme_part.blob)
    return el


def copy_theme_colors(user_master, base_master):
    try:
        u_theme = user_master.part.part_related_by(THEME_RTYPE)
        b_theme = base_master.part.part_related_by(THEME_RTYPE)
    except Exception:
        return
    uscheme = _theme_el(u_theme).find(aq('clrScheme'))
    bscheme = _theme_el(b_theme).find(aq('clrScheme'))
    if uscheme is None or bscheme is None:
        return
    parent = bscheme.getparent()
    parent.replace(bscheme, copy.deepcopy(uscheme))


# ---------------------------------------------------------------------------
# 6) 按亮度设置文字色（只改 文本 defRPr/rPr 里的 srgbClr，不动形状填充）
# ---------------------------------------------------------------------------
def set_text_color_by_luminance(base_prs, kind, data):
    color = bg_luminance_color(kind, data)
    if color == 'FFFFFF':
        return  # 与内置默认白字一致，无需改动
    m = base_prs.slide_masters[0]
    targets = [m] + list(m.slide_layouts)
    # 文本色 srgbClr 的结构：a:solidFill -> a:defRPr / a:rPr / a:endParaRPr
    # 直接父节点是 solidFill，祖父节点才是 run 属性。
    # 形状填充(spPr/ln/sp3d)与 hiddenFill 扩展不在文本语境内，必须排除。
    text_run_props = (aq('defRPr'), aq('rPr'), aq('endParaRPr'))
    cnt = 0
    for part in targets:
        for el in part.element.iter(aq('srgbClr')):
            if el.get('val') != 'FFFFFF':
                continue  # 只翻转白字（在浅底上不可见），保留既有强调色
            p = el.getparent()              # 期望 a:solidFill
            if p is None or p.tag != aq('solidFill'):
                continue
            gp = p.getparent()              # 期望 defRPr / rPr / endParaRPr
            if gp is None or gp.tag not in text_run_props:
                continue
            el.set('val', color)
            cnt += 1
    print('  [亮度选色] 背景偏%s，正文文字设为 #%s（改动 %d 处）'
          % ('浅' if color == '000000' else '深', color, cnt))


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------
def restyle(user_path=None, image_path=None, base_path=None, output=None):
    if base_path is None:
        base_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                 '..', 'assets', '轻风模板_规范版.pptx')
    base_path = os.path.abspath(base_path)
    if output is None:
        output = '更新后模板.pptx'

    user_prs = Presentation(user_path) if user_path else None

    # —— 情况 A：上传了规范模板 → 直接用，不换肤 ——
    if user_prs is not None and is_standard_template(user_prs):
        import shutil
        shutil.copyfile(user_path, output)
        print('[检测结果] 上传模板为规范母版（含 标题/关键词/正文 三级占位符），直接使用，不换肤。')
        print('已生成 ->', output)
        return

    # —— 准备背景来源 ——
    # set_master_bg 接收：image -> (blob_bytes, ext)；solid -> 'RRGGBB'
    bg_kind = None
    bg_data = None
    if image_path:
        ext = os.path.splitext(image_path)[1].lstrip('.') or 'png'
        with open(image_path, 'rb') as f:
            bg_data = (f.read(), ext)
        bg_kind = 'image'
    elif user_prs is not None:
        bg = get_master_bg(user_prs)
        if bg is None:
            # 母版无背景，且未给图片
            print('上传模板非规范母版，请上传一张图片作为生成PPT背景')
            sys.exit(2)
        if bg[0] == 'image':
            bg_kind = 'image'
            bg_data = (bg[1], bg[2])
        else:
            bg_kind = 'solid'
            bg_data = bg[1]

    # —— 以轻风骨架为基，换肤 ——
    base = Presentation(base_path)
    base_master = base.slide_masters[0]
    info = set_master_bg(base_master, bg_kind, bg_data,
                         base.slide_width, base.slide_height)  # ('image',bytes) / ('solid',color)
    if user_prs is not None:
        copy_theme_colors(user_prs.slide_masters[0], base_master)
    set_text_color_by_luminance(base, info[0], info[1])
    base.save(output)
    print('已生成（换肤后）->', output)


if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('--user', help='用户上传的模板 pptx')
    ap.add_argument('--image', help='背景图片（优先作背景）')
    ap.add_argument('--base', help='内置骨架模板（默认 轻风模板_规范版.pptx）')
    ap.add_argument('--output', '-o', required=True, help='输出路径')
    args = ap.parse_args()
    if not args.user and not args.image:
        print('至少需要 --user 或 --image 之一')
        sys.exit(1)
    restyle(user_path=args.user, image_path=args.image,
            base_path=args.base, output=args.output)
